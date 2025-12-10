#!/usr/bin/env python3
import argparse
import pdfplumber
from pptx import Presentation
import sys
import os


def pdf_to_ppt(pdf_path, ppt_path):
    """将 PDF 转换为 PowerPoint"""
    try:
        prs = Presentation()

        with pdfplumber.open(pdf_path) as pdf:
            for i, page in enumerate(pdf.pages):
                # 创建新幻灯片
                slide_layout = prs.slide_layouts[1]  # 标题和内容布局
                slide = prs.slides.add_slide(slide_layout)

                # 添加标题
                title = slide.shapes.title
                title.text = f"第 {i+1} 页"

                # 添加内容
                content = slide.placeholders[1]
                text_frame = content.text_frame
                text_frame.text = page.extract_text() or "无文本内容"

        prs.save(ppt_path)
        print(f"转换成功: {pdf_path} -> {ppt_path}")
        return True
    except Exception as e:
        print(f"转换失败: {str(e)}")
        return False


def main():
    parser = argparse.ArgumentParser(description="PDF 转 PowerPoint")
    parser.add_argument("-i", "--input", required=True, help="输入 PDF 文件路径")
    parser.add_argument("-o", "--output", required=True, help="输出 PowerPoint 文件路径")

    args = parser.parse_args()

    if not os.path.exists(args.input):
        print(f"错误: 输入文件不存在 {args.input}")
        sys.exit(1)

    success = pdf_to_ppt(args.input, args.output)
    sys.exit(0 if success else 1)


if __name__ == "__main__":
    main()
